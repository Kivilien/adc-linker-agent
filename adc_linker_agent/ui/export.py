"""
ADC Linker 设计报告导出模块。

支持三种格式:
  - HTML 幻灯片 (主推): Jinja2 模板 + Nord 主题，浏览器直接演示
  - PDF: reportlab 多页文档
  - PPTX: python-pptx 可编辑演示文稿

所有函数均接受 DesignReport 对象，返回 bytes | None。
库未安装时返回 None，调用方负责降级提示。
"""

from __future__ import annotations

import base64
import io
import textwrap
from pathlib import Path
from typing import Any

# ═══════════════════════════════════════════════════════════════
# 工具函数
# ═══════════════════════════════════════════════════════════════


def _get_2d_image_b64(smiles: str, size: tuple[int, int] = (400, 250)) -> str:
    """生成 2D 分子结构图的 base64 data URI。"""
    from adc_linker_agent.domain.molecule import render_molecule_image

    png_bytes = render_molecule_image(smiles, size=size)
    if png_bytes:
        b64 = base64.b64encode(png_bytes).decode("ascii")
        return f"data:image/png;base64,{b64}"
    return ""


def _get_2d_image_bytes(smiles: str, size: tuple[int, int] = (300, 180)) -> bytes | None:
    """生成 2D 分子结构图 PNG bytes。"""
    from adc_linker_agent.domain.molecule import render_molecule_image

    return render_molecule_image(smiles, size=size)


def _score_color(score: float) -> str:
    """分数 → Nord 色值。"""
    if score >= 0.8:
        return "var(--success)"
    elif score >= 0.5:
        return "var(--accent)"
    elif score >= 0.3:
        return "var(--warning)"
    return "var(--danger)"


# ═══════════════════════════════════════════════════════════════
# HTML 幻灯片导出（主推格式）
# ═══════════════════════════════════════════════════════════════

_TEMPLATE_PATH = Path(__file__).parent / "templates" / "slides.html"


def generate_html_slides(report: Any) -> bytes | None:
    """
    生成 Nord 主题 HTML 幻灯片。

    参考 html-ppt-skill (lewislulu) 和 guizang-ppt-skill 的设计理念:
      - 单文件 HTML，浏览器直接打开演示
      - 键盘/触摸导航
      - Nord 配色，与项目视觉一致

    Args:
        report: domain.report.DesignReport 实例

    Returns:
        HTML bytes，或 None（Jinja2 未安装）
    """
    try:
        from jinja2 import Template
    except ImportError:
        return None

    if report is None or not hasattr(report, "candidates"):
        return None

    template_str = _TEMPLATE_PATH.read_text(encoding="utf-8")
    template = Template(template_str)

    slides = [
        _build_title_slide(report),
        _build_overview_slide(report),
        _build_table_slide(report),
    ]

    for card in (report.detailed_cards or [])[:3]:
        slides.append(_build_candidate_slide(card))

    slides.append(_build_toxicity_slide(report))

    html = template.render(
        report_title=report.request_summary or "ADC Linker Design",
        slides=slides,
    )
    return html.encode("utf-8")


def _build_title_slide(report: Any) -> str:
    """封面幻灯片: 标题 + 参数摘要。"""
    return textwrap.dedent(f"""\
    <div style="display:flex;flex-direction:column;justify-content:center;align-items:center;text-align:center;height:100%;">
      <div style="font-size:0.85em;color:var(--nord9);letter-spacing:4px;text-transform:uppercase;margin-bottom:20px;">ADC Linker Intelligence</div>
      <h1 style="font-size:3.5em;margin-bottom:16px;">ADC 连接子设计报告</h1>
      <p class="subtitle">{report.request_summary}</p>
      <div style="margin-top:30px;">
        <span class="tag tag-info">生成时间: {report.generated_at}</span>
      </div>
      <div style="margin-top:40px;color:var(--muted);font-size:0.85em;">
        <p>ADC Linker Agent v1.1.0 · AI-Assisted Design</p>
        <p style="margin-top:4px;">本报告仅供科研参考，不可用于临床决策</p>
      </div>
    </div>
    <div class="footer">
      <span>ADC Linker Agent</span>
      <span>1 / {_slide_count(report)}</span>
    </div>""")


def _build_overview_slide(report: Any) -> str:
    """概览幻灯片: 4 个指标卡片。"""
    tox_label = "⚠ 有警报" if report.has_any_toxicity else "✓ 通过"
    tox_class = "danger" if report.has_any_toxicity else "success"
    return textwrap.dedent(f"""\
    <h2>设计概览</h2>
    <div class="metric-grid">
      <div class="metric-card">
        <div class="metric-value">{report.total_evaluated}</div>
        <div class="metric-label">评估候选</div>
      </div>
      <div class="metric-card">
        <div class="metric-value">{report.total_filtered}</div>
        <div class="metric-label">过滤排除</div>
      </div>
      <div class="metric-card">
        <div class="metric-value">{report.candidate_count}</div>
        <div class="metric-label">最终候选</div>
      </div>
      <div class="metric-card">
        <div class="metric-value {tox_class}">{tox_label}</div>
        <div class="metric-label">毒性筛查</div>
      </div>
    </div>
    <p class="subtitle" style="margin-top:10px;">需求: {report.request_summary}</p>
    <div class="footer">
      <span>ADC Linker Agent</span>
      <span>2 / {_slide_count(report)}</span>
    </div>""")


def _build_table_slide(report: Any) -> str:
    """候选对比表幻灯片。"""
    rows = ""
    for c in report.candidates or []:
        blood = "✓" if c.blood_stable else "✗"
        lyso = "✓" if c.lysosome_labile else "—"
        tox = f"⚠ {c.toxicity_count}" if c.has_toxicity_alerts else "✓"
        rows += f"""\
        <tr>
          <td><strong>#{c.rank}</strong></td>
          <td>{c.name}</td>
          <td><span class="tag tag-info">{c.mechanism_label}</span></td>
          <td><strong>{c.overall_score:.3f}</strong></td>
          <td>{blood}</td>
          <td>{lyso}</td>
          <td>{c.qed:.3f}</td>
          <td>{c.logp}</td>
          <td>{c.sas}</td>
          <td>{tox}</td>
        </tr>"""

    return textwrap.dedent(f"""\
    <h2>候选对比表</h2>
    <div style="overflow-x:auto;">
    <table>
      <thead>
        <tr>
          <th>排名</th><th>名称</th><th>机制</th><th>综合分</th>
          <th>血液</th><th>溶酶体</th><th>QED</th><th>LogP</th><th>SAS</th><th>毒性</th>
        </tr>
      </thead>
      <tbody>{rows}
      </tbody>
    </table>
    </div>
    <div class="footer">
      <span>ADC Linker Agent</span>
      <span>3 / {_slide_count(report)}</span>
    </div>""")


def _build_candidate_slide(card: dict) -> str:
    """单个候选详细卡片幻灯片。"""
    rank = card.get("rank", "?")
    name = card.get("name", "")
    smiles = card.get("smiles", "")
    mechanism = card.get("mechanism_label", card.get("mechanism", ""))
    scores = card.get("scores", {})
    strengths = card.get("strengths", [])
    weaknesses = card.get("weaknesses", [])
    rec = card.get("recommendation", "")
    ph = card.get("ph_stability", {})

    # 分子图
    mol_img_tag = ""
    b64 = _get_2d_image_b64(smiles)
    if b64:
        mol_img_tag = f'<img src="{b64}" alt="{name}" class="mol-img" style="max-height:200px;">'

    # 分数条
    score_bars = ""
    score_fields = [
        ("血液稳定性", scores.get("blood_stability", 0)),
        ("溶酶体裂解", scores.get("lysosome_lability", 0)),
        ("药物相似性", scores.get("drug_likeness", 0)),
        ("合成可行性", scores.get("synthetic", 0)),
        ("综合分", scores.get("overall", 0)),
    ]
    for label, val in score_fields:
        pct = int(val * 100)
        color = _score_color(val)
        score_bars += f"""\
        <div class="score-item">
          <div class="score-label">{label}</div>
          <div class="score-bar-bg">
            <div class="score-bar-fill" style="width:{pct}%;background:{color};"></div>
          </div>
          <div class="score-value">{val:.2f}</div>
        </div>"""

    # pH 稳定性
    blood_stable = ph.get("blood_stable", False)
    lyso_labile = ph.get("lysosome_labile", False)
    ph_html = f"""\
    <div style="margin-top:12px;">
      <span class="tag {"tag-ok" if blood_stable else "tag-danger"}">
        {"✓ 血液 pH 7.4 稳定" if blood_stable else "✗ 血液 pH 7.4 不稳定"}
      </span>
      <span class="tag {"tag-ok" if lyso_labile else "tag-warn"}">
        {"✓ 溶酶体 pH 5.0 可裂解" if lyso_labile else "— 溶酶体裂解不充分"}
      </span>
    </div>"""

    # 优缺点
    sw_html = ""
    if strengths or weaknesses:
        sw_html = '<div class="cols-2" style="margin-top:12px;"><div>'
        if strengths:
            sw_html += '<h3>✅ 优势</h3><ul class="sw-list">'
            for s in strengths:
                sw_html += f"<li>{s}</li>"
            sw_html += "</ul>"
        sw_html += "</div><div>"
        if weaknesses:
            sw_html += '<h3>⚠️ 不足</h3><ul class="sw-list">'
            for w in weaknesses:
                sw_html += f"<li>{w}</li>"
            sw_html += "</ul>"
        sw_html += "</div></div>"

    # 推荐理由
    rec_html = ""
    if rec:
        rec_class = "tag-danger" if "🚨" in rec else ("tag-ok" if "✅" in rec else "tag-warn")
        rec_html = f'<div style="margin-top:12px;"><span class="tag {rec_class}">{rec}</span></div>'

    rank_emoji = {1: "🥇", 2: "🥈", 3: "🥉"}.get(rank, f"#{rank}")

    return textwrap.dedent(f"""\
    <h2>{rank_emoji} #{rank} — {name}</h2>
    <p style="margin-bottom:4px;"><span class="tag tag-info">{mechanism}</span></p>
    <p class="smiles-mono">{smiles}</p>

    <div class="cols-2" style="margin-top:16px;">
      <div>{mol_img_tag}</div>
      <div>
        <h3>评分</h3>
        <div class="score-row">{score_bars}</div>
        {ph_html}
      </div>
    </div>

    {sw_html}
    {rec_html}

    <div class="footer">
      <span>ADC Linker Agent</span>
      <span>{rank} / Top Candidate</span>
    </div>""")


def _build_toxicity_slide(report: Any) -> str:
    """毒性评估 + 警告幻灯片。"""
    tox_status = "⚠ 检测到毒性警报" if report.has_any_toxicity else "✓ 未检出已知毒性结构"
    tox_class = "danger" if report.has_any_toxicity else "success"
    warnings_html = ""
    if report.warnings:
        warnings_html = '<div style="margin-top:20px;"><h3>⚠️ 警告</h3><ul class="sw-list">'
        for w in report.warnings:
            warnings_html += f'<li style="color:var(--warning);">{w}</li>'
        warnings_html += "</ul></div>"

    failed_html = ""
    if report.failed_scaffolds:
        failed_html = '<div style="margin-top:20px;"><h3>评估失败的骨架</h3><ul class="sw-list">'
        for f in report.failed_scaffolds[:10]:
            failed_html += f"<li>{f.get('name', '?')}: {f.get('error', 'Unknown error')}</li>"
        failed_html += "</ul></div>"

    return textwrap.dedent(f"""\
    <h2>安全性评估</h2>
    <div class="metric-grid">
      <div class="metric-card" style="grid-column:span 2;">
        <div class="metric-value {tox_class}">{tox_status}</div>
        <div class="metric-label">毒性筛查结果</div>
      </div>
    </div>
    <p style="margin-top:16px;">{report.toxicity_summary}</p>
    {warnings_html}
    {failed_html}
    <div style="margin-top:40px;padding:20px;background:var(--card-bg);border-radius:10px;border:1px solid var(--border);">
      <p style="color:var(--muted);font-size:0.85em;">
        ⚕ 本报告由 AI 辅助生成，仅供科研参考，不可用于临床决策。所有连接子候选均需经过实验验证。
      </p>
    </div>
    <div class="footer">
      <span>ADC Linker Agent</span>
      <span>最后</span>
    </div>""")


def _slide_count(report: Any) -> int:
    """计算幻灯片总数。"""
    cards = len(report.detailed_cards) if report.detailed_cards else 0
    return 3 + min(cards, 3) + 1  # 标题+概览+表格+N张卡片+毒性


# ═══════════════════════════════════════════════════════════════
# PPTX 导出
# ═══════════════════════════════════════════════════════════════


def generate_pptx(report: Any) -> bytes | None:
    """
    生成可编辑 PPTX 演示文稿。

    结构: 标题幻灯片 → 候选对比表 → 详细卡片 → 毒性汇总
    使用 Nord 配色方案。

    Args:
        report: domain.report.DesignReport 实例

    Returns:
        PPTX bytes，或 None（python-pptx 未安装）
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except ImportError:
        return None

    if report is None:
        return None

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Nord colors
    accent = RGBColor(0x88, 0xC0, 0xD0)
    dark = RGBColor(0x2E, 0x34, 0x40)
    white = RGBColor(0xEC, 0xEF, 0xF4)
    gray = RGBColor(0x4C, 0x56, 0x6A)
    green = RGBColor(0xA3, 0xBE, 0x8C)
    red = RGBColor(0xBF, 0x61, 0x6A)

    def _dark_slide(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = dark

    def _add_title(slide, text, left=1, top=0.5, width=11, height=0.8, size=28):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = accent
        p.font.bold = True
        return tf

    def _add_text(slide, text, left, top, width, height, size=14, color=None):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color or white
        return tf

    # ─── Slide 1: Title ───
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide(sl)
    _add_title(sl, "ADC Linker Design Report", top=1.5, size=40)
    _add_text(sl, f"Generated: {report.generated_at}", 1, 3, 11, 0.5, size=16, color=gray)
    _add_text(sl, f"Request: {report.request_summary}", 1, 3.8, 11, 0.5, size=16)
    _add_text(
        sl,
        (
            f"Candidates: {report.total_evaluated} evaluated, "
            f"{report.total_filtered} filtered, "
            f"{report.candidate_count} final | "
            f"Toxicity: {'Alerts' if report.has_any_toxicity else 'Pass'}"
        ),
        1,
        5,
        11,
        1,
        size=14,
        color=gray,
    )

    # ─── Slide 2: Candidate Table ───
    sl2 = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide(sl2)
    _add_title(sl2, "Candidate Comparison")

    if report.candidates:
        rows = len(report.candidates) + 1
        cols = 7
        tbl_shape = sl2.shapes.add_table(rows, cols, Inches(1), Inches(1.8), Inches(11), Inches(4))
        tbl = tbl_shape.table

        for i, w in enumerate([0.6, 2.2, 2.2, 1, 1, 1, 1]):
            tbl.columns[i].width = Inches(w)

        headers = ["Rank", "Name", "Mechanism", "Score", "QED", "LogP", "SAS"]
        for i, h in enumerate(headers):
            cell = tbl.cell(0, i)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = dark
                p.font.bold = True

        for ri, c in enumerate(report.candidates, 1):
            vals = [
                str(c.rank),
                c.name,
                c.mechanism,
                f"{c.overall_score:.3f}",
                f"{c.qed:.3f}",
                str(c.logp),
                str(c.sas),
            ]
            for ci, v in enumerate(vals):
                cell = tbl.cell(ri, ci)
                cell.text = v
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.color.rgb = white

    # ─── Slides 3+: Candidate Detail Cards ───
    for card in (report.detailed_cards or [])[:3]:
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        _dark_slide(sl)

        name = card.get("name", "")
        rank = card.get("rank", "?")
        smiles = card.get("smiles", "")
        mechanism = card.get("mechanism_label", card.get("mechanism", ""))
        scores = card.get("scores", {})
        strengths = card.get("strengths", [])
        weaknesses = card.get("weaknesses", [])

        _add_title(sl, f"#{rank} — {name}")
        _add_text(sl, f"Mechanism: {mechanism}", 1, 1.3, 11, 0.4, size=13, color=accent)
        _add_text(sl, f"SMILES: {smiles}", 1, 1.7, 11, 0.4, size=9, color=gray)

        # 2D image
        img_bytes = _get_2d_image_bytes(smiles)
        if img_bytes:
            from io import BytesIO

            sl.shapes.add_picture(
                BytesIO(img_bytes), Inches(1), Inches(2.5), Inches(4), Inches(2.5)
            )

        # Scores
        score_items = [
            ("Blood Stability", scores.get("blood_stability", 0)),
            ("Lysosome Lability", scores.get("lysosome_lability", 0)),
            ("Drug-likeness", scores.get("drug_likeness", 0)),
            ("Synthetic", scores.get("synthetic", 0)),
            ("Overall", scores.get("overall", 0)),
        ]
        score_text = "  |  ".join(f"{label}: {v:.2f}" for label, v in score_items)
        _add_text(sl, "Scores:", 6, 2.5, 6, 0.3, size=14, color=accent)
        _add_text(sl, score_text, 6, 2.9, 6, 0.5, size=12)

        # Strengths/Weaknesses
        if strengths:
            _add_text(
                sl, f"Strengths: {'; '.join(strengths[:3])}", 6, 3.8, 6, 1, size=11, color=green
            )
        if weaknesses:
            _add_text(
                sl, f"Weaknesses: {'; '.join(weaknesses[:3])}", 6, 4.8, 6, 1, size=11, color=red
            )

        # Recommendation
        rec = card.get("recommendation", "")
        if rec:
            _add_text(sl, rec, 1, 5.5, 11, 0.8, size=12, color=accent)

    # ─── Final: Toxicity ───
    sl_final = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide(sl_final)
    _add_title(sl_final, "Safety Assessment")
    _add_text(sl_final, report.toxicity_summary, 1, 2, 11, 3, size=15)

    if report.warnings:
        _add_text(sl_final, "Warnings:", 1, 4.5, 11, 0.5, size=16, color=red)
        for wi, w in enumerate(report.warnings):
            _add_text(sl_final, f"• {w}", 1.5, 5 + wi * 0.4, 10, 0.4, size=12, color=red)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════
# PDF 导出
# ═══════════════════════════════════════════════════════════════


def generate_pdf(report: Any) -> bytes | None:
    """
    生成 Nord 主题多页 PDF 报告。

    结构: 标题页 → 概览指标 → 候选对比表 → Top-3 详细卡片 → 毒性评估

    Args:
        report: domain.report.DesignReport 实例

    Returns:
        PDF bytes，或 None（reportlab 未安装）
    """
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import mm
        from reportlab.platypus import (
            PageBreak,
            Paragraph,
            SimpleDocTemplate,
            Spacer,
            Table,
            TableStyle,
        )
    except ImportError:
        return None

    if report is None:
        return None

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        rightMargin=20 * mm,
        leftMargin=20 * mm,
        topMargin=20 * mm,
        bottomMargin=20 * mm,
        title="ADC Linker Design Report",
        author="ADC Linker Agent",
    )

    # ─── Nord 配色 ───
    nord_dark = colors.HexColor("#2E3440")
    nord_light = colors.HexColor("#ECEFF4")
    nord_accent = colors.HexColor("#88C0D0")
    nord_gray = colors.HexColor("#4C566A")
    nord_green = colors.HexColor("#A3BE8C")
    nord_red = colors.HexColor("#BF616A")
    nord_yellow = colors.HexColor("#EBCB8B")

    # ─── 样式 ───
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "NordTitle",
        parent=styles["Title"],
        fontSize=22,
        textColor=nord_accent,
        spaceAfter=12,
    )
    h2_style = ParagraphStyle(
        "NordH2",
        parent=styles["Heading2"],
        fontSize=16,
        textColor=nord_accent,
        spaceAfter=8,
    )
    h3_style = ParagraphStyle(
        "NordH3",
        parent=styles["Heading3"],
        fontSize=13,
        textColor=nord_accent,
        spaceAfter=6,
    )
    body_style = ParagraphStyle(
        "NordBody",
        parent=styles["Normal"],
        fontSize=10,
        textColor=nord_dark,
        leading=14,
    )
    small_style = ParagraphStyle(
        "NordSmall",
        parent=styles["Normal"],
        fontSize=8,
        textColor=nord_gray,
    )
    mono_style = ParagraphStyle(
        "NordMono",
        parent=styles["Normal"],
        fontSize=7,
        textColor=nord_gray,
        fontName="Courier",
    )

    story = []

    def _add_spacer(height=6):
        story.append(Spacer(1, height * mm))

    # ─── 第 1 页: 标题 ───
    story.append(Paragraph("ADC 连接子设计报告", title_style))
    _add_spacer(4)
    story.append(Paragraph(f"需求: {report.request_summary or 'N/A'}", body_style))
    story.append(Paragraph(f"生成时间: {report.generated_at or 'N/A'}", small_style))
    _add_spacer(10)

    # 概览指标表格
    overview_data = [
        ["评估候选", "过滤排除", "最终候选", "毒性筛查"],
        [
            str(report.total_evaluated),
            str(report.total_filtered),
            str(report.candidate_count),
            "⚠ 警报" if report.has_any_toxicity else "✓ 通过",
        ],
    ]
    overview_table = Table(overview_data, colWidths=[60, 60, 60, 80])
    overview_table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), nord_accent),
                ("TEXTCOLOR", (0, 0), (-1, 0), nord_light),
                ("BACKGROUND", (0, 1), (-1, 1), nord_light),
                ("TEXTCOLOR", (0, 1), (-1, 1), nord_dark),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("FONTSIZE", (0, 0), (-1, -1), 10),
                ("GRID", (0, 0), (-1, -1), 0.5, nord_gray),
                ("ROWBACKGROUNDS", (0, 1), (-1, 1), [nord_light]),
            ]
        )
    )
    story.append(overview_table)
    story.append(PageBreak())

    # ─── 第 2 页: 候选对比表 ───
    story.append(Paragraph("候选对比表", h2_style))
    _add_spacer(4)

    if report.candidates:
        headers = ["排名", "名称", "机制", "综合分", "QED", "LogP", "SAS", "毒性"]
        table_data = [headers]
        for c in report.candidates:
            tox = f"⚠ {c.toxicity_count}" if c.has_toxicity_alerts else "✓"
            table_data.append(
                [
                    str(c.rank),
                    c.name,
                    c.mechanism_label or c.mechanism,
                    f"{c.overall_score:.3f}",
                    f"{c.qed:.3f}",
                    str(c.logp),
                    str(c.sas),
                    tox,
                ]
            )

        col_widths = [30, 70, 65, 45, 35, 30, 30, 35]
        cand_table = Table(table_data, colWidths=col_widths)
        cand_table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), nord_dark),
                    ("TEXTCOLOR", (0, 0), (-1, 0), nord_light),
                    ("FONTSIZE", (0, 0), (-1, -1), 7),
                    ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                    ("ALIGN", (1, 0), (1, -1), "LEFT"),
                    ("GRID", (0, 0), (-1, -1), 0.3, nord_gray),
                ]
            )
        )
        story.append(cand_table)
    else:
        story.append(Paragraph("无候选数据", body_style))

    story.append(PageBreak())

    # ─── 第 3-5 页: Top-3 详细卡片 ───
    for card in (report.detailed_cards or [])[:3]:
        rank = card.get("rank", "?")
        name = card.get("name", "")
        smiles = card.get("smiles", "")
        mechanism = card.get("mechanism_label", card.get("mechanism", ""))
        scores = card.get("scores", {})
        strengths = card.get("strengths", [])
        weaknesses = card.get("weaknesses", [])
        rec = card.get("recommendation", "")

        rank_emoji = {1: "🥇", 2: "🥈", 3: "🥉"}.get(rank, f"#{rank}")
        story.append(Paragraph(f"{rank_emoji} #{rank} — {name}", h2_style))
        story.append(Paragraph(f"机制: {mechanism}", body_style))
        story.append(Paragraph(f"SMILES: {smiles}", mono_style))
        _add_spacer(4)

        # 分子结构图
        img_bytes = _get_2d_image_bytes(smiles)
        if img_bytes:
            from io import BytesIO as Bio

            from reportlab.platypus import Image

            img = Image(Bio(img_bytes), width=120, height=75)
            story.append(img)
            _add_spacer(4)

        # 评分
        score_items = [
            ("血液稳定性", scores.get("blood_stability", 0)),
            ("溶酶体裂解", scores.get("lysosome_lability", 0)),
            ("药物相似性", scores.get("drug_likeness", 0)),
            ("合成可行性", scores.get("synthetic", 0)),
            ("综合分", scores.get("overall", 0)),
        ]
        score_text = " | ".join(f"{label}: {val:.2f}" for label, val in score_items)
        story.append(Paragraph(f"评分: {score_text}", body_style))
        _add_spacer(4)

        if strengths:
            story.append(Paragraph(f"✅ 优势: {'; '.join(strengths[:3])}", body_style))
        if weaknesses:
            story.append(Paragraph(f"⚠️ 不足: {'; '.join(weaknesses[:3])}", body_style))
        if rec:
            story.append(Paragraph(f"💡 推荐: {rec}", body_style))

        story.append(PageBreak())

    # ─── 最后: 毒性评估 ───
    story.append(Paragraph("安全性评估", h2_style))
    _add_spacer(4)
    tox_status = "⚠ 检测到毒性警报" if report.has_any_toxicity else "✓ 未检出已知毒性结构"
    tox_color = nord_red if report.has_any_toxicity else nord_green
    story.append(
        Paragraph(
            tox_status,
            ParagraphStyle("ToxStatus", parent=body_style, textColor=tox_color, fontSize=12),
        )
    )
    _add_spacer(4)

    if report.toxicity_summary:
        story.append(Paragraph(report.toxicity_summary, body_style))

    if report.warnings:
        _add_spacer(4)
        story.append(Paragraph("⚠️ 警告:", h3_style))
        for w in report.warnings:
            story.append(
                Paragraph(
                    f"• {w}", ParagraphStyle("WarnItem", parent=body_style, textColor=nord_yellow)
                )
            )

    _add_spacer(10)
    story.append(Paragraph("⚕ 本报告由 AI 辅助生成，仅供科研参考，不可用于临床决策。", small_style))
    story.append(Paragraph(f"ADC Linker Agent v1.1.0 — {report.generated_at or ''}", small_style))

    doc.build(story)
    return buf.getvalue()
